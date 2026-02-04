"""
Enterprise Document Processing Engine
Handles PDF, PPT, DOCX, Images, and other business document formats
"""

import pandas as pd
import numpy as np
from typing import Dict, Any, List, Optional, Tuple, Union
from dataclasses import dataclass
from enum import Enum
import os
import io
import base64
from datetime import datetime
import warnings
warnings.filterwarnings('ignore')

# Document processing imports
try:
    import PyPDF2
    import pdfplumber
    from docx import Document
    from pptx import Presentation
    import tabula
    # import camelot  # Uncomment when needed
    from PIL import Image
    # import pytesseract  # Uncomment when OCR is needed
except ImportError as e:
    print(f"Some document processing libraries not available: {e}")

# Text processing
try:
    import nltk
    from textblob import TextBlob
    from wordcloud import WordCloud
except ImportError:
    print("Text processing libraries not fully available")

class DocumentType(Enum):
    PDF = "pdf"
    DOCX = "docx"
    DOC = "doc"
    PPTX = "pptx"
    PPT = "ppt"
    TXT = "txt"
    RTF = "rtf"
    IMAGE = "image"
    HTML = "html"
    EMAIL = "email"
    UNKNOWN = "unknown"

class ExtractionMethod(Enum):
    TEXT_EXTRACTION = "text_extraction"
    TABLE_EXTRACTION = "table_extraction"
    OCR = "ocr"
    METADATA_EXTRACTION = "metadata"
    HYBRID = "hybrid"

@dataclass
class DocumentMetadata:
    """Comprehensive document metadata"""
    filename: str
    document_type: DocumentType
    file_size_bytes: int
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    character_count: Optional[int] = None
    language: Optional[str] = None
    
    # Document properties
    title: Optional[str] = None
    author: Optional[str] = None
    subject: Optional[str] = None
    creator: Optional[str] = None
    creation_date: Optional[str] = None
    modification_date: Optional[str] = None
    
    # Content analysis
    has_tables: bool = False
    has_images: bool = False
    table_count: int = 0
    image_count: int = 0
    
    # Processing info
    extraction_method: ExtractionMethod = ExtractionMethod.TEXT_EXTRACTION
    processing_time: float = 0.0
    confidence_score: float = 0.0
    warnings: List[str] = None

@dataclass
class ProcessedDocument:
    """Processed document with extracted data"""
    metadata: DocumentMetadata
    
    # Extracted content
    text_content: Optional[str] = None
    structured_data: Optional[pd.DataFrame] = None
    tables: Optional[List[pd.DataFrame]] = None
    
    # Text analytics
    sentiment_score: Optional[float] = None
    key_phrases: Optional[List[str]] = None
    entities: Optional[List[Dict]] = None
    summary: Optional[str] = None
    
    # Quality indicators
    extraction_quality: str = "Unknown"
    data_completeness: float = 0.0
    
    is_valid: bool = True
    error_message: Optional[str] = None

class EnterpriseDocumentProcessor:
    """
    Enterprise-grade document processing engine
    Handles all major business document formats
    """
    
    def __init__(self):
        self.supported_extensions = {
            '.pdf': DocumentType.PDF,
            '.docx': DocumentType.DOCX,
            '.doc': DocumentType.DOC,
            '.pptx': DocumentType.PPTX,
            '.ppt': DocumentType.PPT,
            '.txt': DocumentType.TXT,
            '.rtf': DocumentType.RTF,
            '.html': DocumentType.HTML,
            '.htm': DocumentType.HTML,
            '.png': DocumentType.IMAGE,
            '.jpg': DocumentType.IMAGE,
            '.jpeg': DocumentType.IMAGE,
            '.tiff': DocumentType.IMAGE,
            '.bmp': DocumentType.IMAGE
        }
        
        # Processing configurations
        self.max_file_size = 100 * 1024 * 1024  # 100MB
        self.ocr_enabled = False  # Enable when tesseract is configured
        
    def process_document(self, file_path: str) -> ProcessedDocument:
        """
        Main document processing method
        """
        start_time = datetime.now()
        
        try:
            # Basic validation
            if not os.path.exists(file_path):
                return self._create_error_result("File not found", file_path)
            
            file_size = os.path.getsize(file_path)
            if file_size > self.max_file_size:
                return self._create_error_result(f"File too large: {file_size} bytes", file_path)
            
            # Detect document type
            doc_type = self._detect_document_type(file_path)
            if doc_type == DocumentType.UNKNOWN:
                return self._create_error_result("Unsupported document type", file_path)
            
            # Initialize metadata
            metadata = DocumentMetadata(
                filename=os.path.basename(file_path),
                document_type=doc_type,
                file_size_bytes=file_size,
                warnings=[]
            )
            
            # Process based on document type
            if doc_type == DocumentType.PDF:
                result = self._process_pdf(file_path, metadata)
            elif doc_type == DocumentType.DOCX:
                result = self._process_docx(file_path, metadata)
            elif doc_type == DocumentType.PPTX:
                result = self._process_pptx(file_path, metadata)
            elif doc_type == DocumentType.TXT:
                result = self._process_text(file_path, metadata)
            elif doc_type == DocumentType.IMAGE:
                result = self._process_image(file_path, metadata)
            else:
                result = self._create_error_result(f"Processing not implemented for {doc_type.value}", file_path)
            
            # Post-processing analytics
            if result.is_valid and result.text_content:
                result = self._add_text_analytics(result)
            
            # Calculate processing time
            processing_time = (datetime.now() - start_time).total_seconds()
            result.metadata.processing_time = processing_time
            
            return result
            
        except Exception as e:
            return self._create_error_result(f"Processing failed: {str(e)}", file_path)
    
    def _detect_document_type(self, file_path: str) -> DocumentType:
        """Detect document type from file extension"""
        _, ext = os.path.splitext(file_path.lower())
        return self.supported_extensions.get(ext, DocumentType.UNKNOWN)
    
    def _process_pdf(self, file_path: str, metadata: DocumentMetadata) -> ProcessedDocument:
        """Process PDF documents"""
        try:
            text_content = ""
            tables = []
            
            # Method 1: Try pdfplumber for better text and table extraction
            try:
                with pdfplumber.open(file_path) as pdf:
                    metadata.page_count = len(pdf.pages)
                    
                    for page in pdf.pages:
                        # Extract text
                        page_text = page.extract_text()
                        if page_text:
                            text_content += page_text + "\n"
                        
                        # Extract tables
                        page_tables = page.extract_tables()
                        if page_tables:
                            for table in page_tables:
                                if table and len(table) > 1:  # Valid table
                                    df = pd.DataFrame(table[1:], columns=table[0])
                                    tables.append(df)
                
                metadata.extraction_method = ExtractionMethod.HYBRID
                
            except Exception as e:
                metadata.warnings.append(f"pdfplumber failed: {str(e)}")
                
                # Fallback: PyPDF2
                try:
                    with open(file_path, 'rb') as file:
                        pdf_reader = PyPDF2.PdfReader(file)
                        metadata.page_count = len(pdf_reader.pages)
                        
                        for page in pdf_reader.pages:
                            text_content += page.extract_text() + "\n"
                    
                    metadata.extraction_method = ExtractionMethod.TEXT_EXTRACTION
                    
                except Exception as e2:
                    return self._create_error_result(f"PDF processing failed: {str(e2)}", file_path)
            
            # Try table extraction with tabula
            try:
                tabula_tables = tabula.read_pdf(file_path, pages='all', multiple_tables=True)
                if tabula_tables:
                    tables.extend(tabula_tables)
                    metadata.has_tables = True
                    metadata.table_count = len(tables)
            except:
                metadata.warnings.append("Tabula table extraction failed")
            
            # Combine tables into structured data
            structured_data = None
            if tables:
                try:
                    # Combine all tables
                    structured_data = pd.concat(tables, ignore_index=True)
                except:
                    structured_data = tables[0] if tables else None
            
            # Update metadata
            metadata.word_count = len(text_content.split()) if text_content else 0
            metadata.character_count = len(text_content) if text_content else 0
            metadata.confidence_score = 0.9 if text_content else 0.3
            
            return ProcessedDocument(
                metadata=metadata,
                text_content=text_content,
                structured_data=structured_data,
                tables=tables,
                extraction_quality="High" if text_content else "Low",
                data_completeness=1.0 if text_content else 0.0,
                is_valid=True
            )
            
        except Exception as e:
            return self._create_error_result(f"PDF processing error: {str(e)}", file_path)
    
    def _process_docx(self, file_path: str, metadata: DocumentMetadata) -> ProcessedDocument:
        """Process DOCX documents"""
        try:
            doc = Document(file_path)
            
            # Extract text
            text_content = ""
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                if table_data and len(table_data) > 1:
                    df = pd.DataFrame(table_data[1:], columns=table_data[0])
                    tables.append(df)
            
            # Extract metadata
            core_props = doc.core_properties
            metadata.title = core_props.title
            metadata.author = core_props.author
            metadata.subject = core_props.subject
            metadata.creation_date = str(core_props.created) if core_props.created else None
            metadata.modification_date = str(core_props.modified) if core_props.modified else None
            
            # Update counts
            metadata.word_count = len(text_content.split()) if text_content else 0
            metadata.character_count = len(text_content) if text_content else 0
            metadata.has_tables = len(tables) > 0
            metadata.table_count = len(tables)
            metadata.confidence_score = 0.95
            metadata.extraction_method = ExtractionMethod.HYBRID
            
            # Combine tables
            structured_data = None
            if tables:
                try:
                    structured_data = pd.concat(tables, ignore_index=True)
                except:
                    structured_data = tables[0]
            
            return ProcessedDocument(
                metadata=metadata,
                text_content=text_content,
                structured_data=structured_data,
                tables=tables,
                extraction_quality="High",
                data_completeness=1.0,
                is_valid=True
            )
            
        except Exception as e:
            return self._create_error_result(f"DOCX processing error: {str(e)}", file_path)
    
    def _process_pptx(self, file_path: str, metadata: DocumentMetadata) -> ProcessedDocument:
        """Process PPTX presentations"""
        try:
            prs = Presentation(file_path)
            
            text_content = ""
            slide_count = 0
            
            for slide in prs.slides:
                slide_count += 1
                slide_text = f"\n--- Slide {slide_count} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                
                text_content += slide_text
            
            # Update metadata
            metadata.page_count = slide_count
            metadata.word_count = len(text_content.split()) if text_content else 0
            metadata.character_count = len(text_content) if text_content else 0
            metadata.confidence_score = 0.9
            metadata.extraction_method = ExtractionMethod.TEXT_EXTRACTION
            
            # Extract presentation properties
            core_props = prs.core_properties
            metadata.title = core_props.title
            metadata.author = core_props.author
            metadata.subject = core_props.subject
            
            return ProcessedDocument(
                metadata=metadata,
                text_content=text_content,
                extraction_quality="High",
                data_completeness=1.0,
                is_valid=True
            )
            
        except Exception as e:
            return self._create_error_result(f"PPTX processing error: {str(e)}", file_path)
    
    def _process_text(self, file_path: str, metadata: DocumentMetadata) -> ProcessedDocument:
        """Process plain text files"""
        try:
            # Detect encoding
            import chardet
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                encoding_result = chardet.detect(raw_data)
                encoding = encoding_result.get('encoding', 'utf-8')
            
            # Read text
            with open(file_path, 'r', encoding=encoding) as f:
                text_content = f.read()
            
            # Update metadata
            metadata.word_count = len(text_content.split()) if text_content else 0
            metadata.character_count = len(text_content) if text_content else 0
            metadata.confidence_score = 1.0
            metadata.extraction_method = ExtractionMethod.TEXT_EXTRACTION
            
            return ProcessedDocument(
                metadata=metadata,
                text_content=text_content,
                extraction_quality="High",
                data_completeness=1.0,
                is_valid=True
            )
            
        except Exception as e:
            return self._create_error_result(f"Text processing error: {str(e)}", file_path)
    
    def _process_image(self, file_path: str, metadata: DocumentMetadata) -> ProcessedDocument:
        """Process image files (with optional OCR)"""
        try:
            # Basic image info
            with Image.open(file_path) as img:
                width, height = img.size
                mode = img.mode
            
            text_content = f"Image: {width}x{height}, Mode: {mode}"
            
            # OCR extraction (if enabled and tesseract available)
            if self.ocr_enabled:
                try:
                    import pytesseract
                    ocr_text = pytesseract.image_to_string(Image.open(file_path))
                    if ocr_text.strip():
                        text_content = ocr_text
                        metadata.extraction_method = ExtractionMethod.OCR
                        metadata.confidence_score = 0.7
                    else:
                        metadata.confidence_score = 0.3
                except:
                    metadata.warnings.append("OCR extraction failed")
                    metadata.confidence_score = 0.2
            else:
                metadata.confidence_score = 0.2
                metadata.warnings.append("OCR not enabled - only basic image info extracted")
            
            metadata.word_count = len(text_content.split()) if text_content else 0
            metadata.character_count = len(text_content) if text_content else 0
            
            return ProcessedDocument(
                metadata=metadata,
                text_content=text_content,
                extraction_quality="Medium" if self.ocr_enabled else "Low",
                data_completeness=0.8 if self.ocr_enabled else 0.2,
                is_valid=True
            )
            
        except Exception as e:
            return self._create_error_result(f"Image processing error: {str(e)}", file_path)
    
    def _add_text_analytics(self, result: ProcessedDocument) -> ProcessedDocument:
        """Add text analytics to processed document"""
        try:
            if not result.text_content:
                return result
            
            text = result.text_content
            
            # Sentiment analysis
            try:
                blob = TextBlob(text)
                result.sentiment_score = blob.sentiment.polarity
            except:
                pass
            
            # Language detection
            try:
                result.metadata.language = blob.detect_language()
            except:
                pass
            
            # Key phrases extraction (simple approach)
            try:
                words = text.lower().split()
                # Remove common stop words
                stop_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'is', 'are', 'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could', 'should'}
                filtered_words = [word for word in words if word not in stop_words and len(word) > 3]
                
                # Get most frequent words
                from collections import Counter
                word_freq = Counter(filtered_words)
                result.key_phrases = [word for word, count in word_freq.most_common(10)]
            except:
                pass
            
            # Simple summary (first few sentences)
            try:
                sentences = text.split('.')[:3]
                result.summary = '. '.join(sentences).strip() + '.'
            except:
                pass
            
            return result
            
        except Exception as e:
            result.metadata.warnings.append(f"Text analytics failed: {str(e)}")
            return result
    
    def _create_error_result(self, error_message: str, file_path: str) -> ProcessedDocument:
        """Create error result"""
        metadata = DocumentMetadata(
            filename=os.path.basename(file_path) if file_path else "unknown",
            document_type=DocumentType.UNKNOWN,
            file_size_bytes=0,
            warnings=[error_message]
        )
        
        return ProcessedDocument(
            metadata=metadata,
            is_valid=False,
            error_message=error_message,
            extraction_quality="Failed",
            data_completeness=0.0
        )
    
    def get_supported_formats(self) -> Dict[str, List[str]]:
        """Get supported document formats"""
        return {
            'documents': ['pdf', 'docx', 'doc', 'pptx', 'ppt', 'txt', 'rtf'],
            'web': ['html', 'htm'],
            'images': ['png', 'jpg', 'jpeg', 'tiff', 'bmp'],
            'max_file_size_mb': self.max_file_size // (1024 * 1024),
            'ocr_enabled': self.ocr_enabled
        }

# Usage example
if __name__ == "__main__":
    processor = EnterpriseDocumentProcessor()
    
    # Test with a sample file
    # result = processor.process_document("sample.pdf")
    # print(f"Processing result: {result.is_valid}")
    # print(f"Text length: {len(result.text_content) if result.text_content else 0}")
    # print(f"Tables found: {len(result.tables) if result.tables else 0}")
    
    print("Enterprise Document Processor initialized")
    print(f"Supported formats: {processor.get_supported_formats()}")